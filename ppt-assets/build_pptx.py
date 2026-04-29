"""
Generates the M4 final-presentation PPTX for EcoRent.

Slide structure mandated by Prof. Hemant Sharma's M4 brief:

  Title  — project + author
  Slide 1 — Problem + Final Problem Statement
  Slide 2 — Key Needs, Functions and Design Priorities
  Slide 3 — Concept Selection and Design Evolution (M3 → M4)
  Slide 4 — Final Embodiment Design (the main slide)
  Slide 5 — Conclusion

16:9 widescreen, EcoRent brand palette.

NOTE on text helpers:

  add_lines(slide, …, lines=[…])  → each item is a NEW PARAGRAPH (new line)
  add_runs(slide,  …, runs=[…])   → all runs are inline in ONE paragraph

Use the right one. Mixing them by mistake stacks textboxes on top of each
other (which is what broke v1).
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# =========================================================
# Brand palette
# =========================================================
INK_900 = RGBColor(0x0F, 0x17, 0x2A)
INK_800 = RGBColor(0x1E, 0x29, 0x3B)
INK_700 = RGBColor(0x33, 0x41, 0x55)
INK_600 = RGBColor(0x47, 0x55, 0x69)
INK_500 = RGBColor(0x64, 0x74, 0x8B)
INK_300 = RGBColor(0xCB, 0xD5, 0xE1)
INK_200 = RGBColor(0xE2, 0xE8, 0xF0)
INK_100 = RGBColor(0xF1, 0xF5, 0xF9)
INK_50  = RGBColor(0xF8, 0xFA, 0xFC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

BRAND_400 = RGBColor(0x4A, 0xDE, 0x80)
BRAND_500 = RGBColor(0x22, 0xC5, 0x5E)
BRAND_600 = RGBColor(0x16, 0xA3, 0x4A)
BRAND_700 = RGBColor(0x15, 0x80, 0x3D)
BRAND_50  = RGBColor(0xEC, 0xFD, 0xF5)
BRAND_100 = RGBColor(0xD1, 0xFA, 0xE5)
BRAND_900 = RGBColor(0x14, 0x53, 0x2D)

AMBER_500 = RGBColor(0xF5, 0x9E, 0x0B)
ROSE_500  = RGBColor(0xF4, 0x3F, 0x5E)
SKY_500   = RGBColor(0x0E, 0xA5, 0xE9)

DARK_PANEL = RGBColor(0x14, 0x1B, 0x2E)

ASSETS = Path(__file__).parent

# =========================================================
# Deck setup
# =========================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# =========================================================
# Shape helpers
# =========================================================
def add_rect(slide, x, y, w, h, fill, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.shadow.inherit = False
    return sh


def add_round(slide, x, y, w, h, fill, *, radius=0.06, line=None,
              line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = radius
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.shadow.inherit = False
    return sh


def add_oval(slide, x, y, w, h, fill, *, alpha=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.shadow.inherit = False
    if alpha is not None:
        _set_alpha(sh, alpha)
    return sh


def _set_alpha(shape, alpha: float):
    sp = shape.fill._xPr  # type: ignore[attr-defined]
    fill = sp.find(qn("a:solidFill"))
    if fill is None:
        return
    srgb = fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    a = etree.SubElement(srgb, qn("a:alpha"))
    a.set("val", str(int(alpha * 100000)))


# =========================================================
# Text helpers
# =========================================================
DEFAULTS = dict(
    size=14, color=INK_700, bold=False, italic=False, font="Inter",
    align=PP_ALIGN.LEFT, line_spacing=1.18, space_before=0,
)


def _make_textbox(slide, x, y, w, h, anchor):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    return tb, tf


def _resolve(item, defaults):
    if isinstance(item, str):
        out = dict(defaults)
        out["text"] = item
        return out
    out = dict(defaults)
    out.update(item)
    out.setdefault("text", "")
    return out


def _apply_run(run, attrs):
    run.text = attrs["text"]
    f = run.font
    f.name = attrs["font"]
    f.size = Pt(attrs["size"])
    f.bold = attrs["bold"]
    f.italic = attrs["italic"]
    f.color.rgb = attrs["color"]


def add_lines(slide, x, y, w, h, lines, *, anchor=MSO_ANCHOR.TOP, **kw):
    """Each entry in `lines` becomes its own paragraph (new line).

    Per-item overrides accepted via dict items: text, size, bold, italic,
    color, align, font, line_spacing, space_before.
    """
    defaults = dict(DEFAULTS, **kw)
    tb, tf = _make_textbox(slide, x, y, w, h, anchor)
    for i, item in enumerate(lines):
        attrs = _resolve(item, defaults)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = attrs["align"]
        p.line_spacing = attrs["line_spacing"]
        if attrs["space_before"]:
            p.space_before = Pt(attrs["space_before"])
        r = p.add_run()
        _apply_run(r, attrs)
    return tb


def add_runs(slide, x, y, w, h, runs, *, anchor=MSO_ANCHOR.TOP, **kw):
    """All entries are inline runs inside ONE paragraph (one line, can wrap)."""
    defaults = dict(DEFAULTS, **kw)
    tb, tf = _make_textbox(slide, x, y, w, h, anchor)
    p = tf.paragraphs[0]
    p.alignment = defaults["align"]
    p.line_spacing = defaults["line_spacing"]
    for item in runs:
        attrs = _resolve(item, defaults)
        r = p.add_run()
        _apply_run(r, attrs)
    return tb


def add_text(slide, x, y, w, h, text, *, anchor=MSO_ANCHOR.TOP, **kw):
    """Single line/paragraph of text — convenience wrapper."""
    return add_lines(slide, x, y, w, h, [text], anchor=anchor, **kw)


# =========================================================
# Slide chrome (header + footer common to every content slide)
# =========================================================
def slide_chrome(slide, *, page_label, page_total, dark=False):
    fg = WHITE if dark else INK_900
    fg_dim = INK_300 if dark else INK_500

    add_rect(slide, Emu(0), Emu(0), SW, Inches(0.06), BRAND_500)

    # Mini logo
    add_round(slide, Inches(0.45), Inches(0.30), Inches(0.32), Inches(0.32),
              BRAND_500, radius=0.25)
    add_text(slide, Inches(0.45), Inches(0.30), Inches(0.32), Inches(0.32),
             {"text": "▲", "color": INK_900, "bold": True, "size": 14,
              "align": PP_ALIGN.CENTER},
             anchor=MSO_ANCHOR.MIDDLE)

    add_runs(slide, Inches(0.85), Inches(0.30), Inches(2.4), Inches(0.32),
             [{"text": "Eco", "color": fg, "bold": True, "size": 13},
              {"text": "Rent", "color": BRAND_600 if not dark else BRAND_400,
               "bold": True, "size": 13}],
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(3.3), Inches(0.30), Inches(7.0), Inches(0.32),
             {"text": "PDD MF-F473  ·  Milestone 4  ·  Final Presentation",
              "size": 10, "color": fg_dim},
             anchor=MSO_ANCHOR.MIDDLE)

    add_text(slide, Inches(11.5), Inches(0.30), Inches(1.4), Inches(0.32),
             {"text": f"Slide {page_label} / {page_total}",
              "size": 10, "color": fg_dim, "align": PP_ALIGN.RIGHT},
             anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# TITLE
# =========================================================
def build_title():
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, INK_900)

    # Soft auroras
    add_oval(s, Inches(-3.5), Inches(-3.5), Inches(11), Inches(11),
             BRAND_700, alpha=0.20)
    add_oval(s, Inches(8.5), Inches(3.0), Inches(7), Inches(7),
             SKY_500, alpha=0.10)

    # ---- Logo row (icon + wordmark side-by-side) ----
    add_round(s, Inches(0.85), Inches(0.85), Inches(0.7), Inches(0.7),
              BRAND_500, radius=0.25)
    add_text(s, Inches(0.85), Inches(0.85), Inches(0.7), Inches(0.7),
             {"text": "▲", "size": 30, "bold": True, "color": INK_900,
              "align": PP_ALIGN.CENTER},
             anchor=MSO_ANCHOR.MIDDLE)
    add_runs(s, Inches(1.70), Inches(0.85), Inches(4.0), Inches(0.7),
             [{"text": "Eco", "size": 28, "bold": True, "color": WHITE},
              {"text": "Rent", "size": 28, "bold": True, "color": BRAND_400}],
             anchor=MSO_ANCHOR.MIDDLE)

    # ---- Pill chip ----
    add_round(s, Inches(0.85), Inches(2.30), Inches(4.05), Inches(0.45),
              BRAND_900, radius=0.5)
    add_text(s, Inches(0.85), Inches(2.30), Inches(4.05), Inches(0.45),
             {"text": "● Trust-First Hyper-Local Marketplace",
              "size": 12, "bold": True, "color": BRAND_100,
              "align": PP_ALIGN.CENTER},
             anchor=MSO_ANCHOR.MIDDLE)

    # ---- Hero — three separate textboxes, no wrapping surprises ----
    add_text(s, Inches(0.85), Inches(2.95), Inches(11.6), Inches(0.95),
             {"text": "Share the adventure.",
              "size": 44, "bold": True, "color": WHITE,
              "line_spacing": 1.05})
    add_text(s, Inches(0.85), Inches(3.95), Inches(11.6), Inches(0.95),
             {"text": "Rent outdoor gear",
              "size": 44, "bold": True, "color": BRAND_400,
              "line_spacing": 1.05})
    add_text(s, Inches(0.85), Inches(4.95), Inches(11.6), Inches(0.95),
             {"text": "from your neighbours.",
              "size": 44, "bold": True,
              "color": RGBColor(0x9C, 0xA3, 0xAF),
              "line_spacing": 1.05})

    # ---- Subtitle ----
    add_text(s, Inches(0.85), Inches(6.05), Inches(11.6), Inches(0.4),
             {"text": "M4 Final Embodiment & Design Refinement",
              "size": 16, "color": INK_300})

    # ---- Live demo pill ----
    add_round(s, Inches(0.85), Inches(6.50), Inches(5.30), Inches(0.40),
              BRAND_900, radius=0.5, line=BRAND_500, line_w=1.0)
    add_runs(s, Inches(0.85), Inches(6.50), Inches(5.30), Inches(0.40),
             [{"text": "🌐  LIVE DEMO  ", "size": 10, "bold": True,
               "color": BRAND_400},
              {"text": "https://pdd-azure.vercel.app",
               "size": 11, "bold": True, "color": WHITE}],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- Author + course ----
    add_text(s, Inches(0.85), Inches(7.00), Inches(11.6), Inches(0.30),
             {"text": "Ashish K. Cheruku  ·  2021B3AB1141P  ·  PDD MF-F473  ·  Submitted to Prof. Hemant Sharma  ·  29 Apr 2026",
              "size": 11, "color": INK_300})


# =========================================================
# SLIDE 1 — Problem + Final Problem Statement
# =========================================================
def build_slide_1(page_total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, INK_50)
    slide_chrome(s, page_label="1", page_total=page_total)

    # ---- Project + team strip (M4 requirement: project title + team on Slide 1) ----
    strip_y = Inches(0.95)
    strip_h = Inches(0.42)
    add_round(s, Inches(0.85), strip_y, Inches(12.0), strip_h, WHITE,
              radius=0.5, line=INK_200)
    add_runs(s, Inches(1.10), strip_y, Inches(7.5), strip_h,
             [{"text": "PROJECT  ", "size": 9, "bold": True, "color": INK_500},
              {"text": "EcoRent — Peer-to-Peer Outdoor Equipment Rental",
               "size": 11, "bold": True, "color": INK_900}],
             anchor=MSO_ANCHOR.MIDDLE)
    add_runs(s, Inches(8.20), strip_y, Inches(4.55), strip_h,
             [{"text": "TEAM  ", "size": 9, "bold": True, "color": INK_500},
              {"text": "Ashish K. Cheruku  ·  2021B3AB1141P",
               "size": 11, "color": INK_700}],
             anchor=MSO_ANCHOR.MIDDLE,
             align=PP_ALIGN.RIGHT)

    add_text(s, Inches(0.85), Inches(1.55), Inches(8), Inches(0.32),
             {"text": "01 — PROBLEM",
              "size": 11, "bold": True, "color": BRAND_600})
    add_text(s, Inches(0.85), Inches(1.88), Inches(11.6), Inches(0.6),
             {"text": "Who we’re solving for, and what hurts.",
              "size": 28, "bold": True, "color": INK_900,
              "line_spacing": 1.1})

    # ---- Two-column body ----
    col_w = Inches(5.85)
    card_y = Inches(2.70)
    card_h = Inches(4.30)
    pad = Inches(0.40)

    # Left: target user
    lx = Inches(0.85)
    add_round(s, lx, card_y, col_w, card_h, WHITE, radius=0.04, line=INK_200)
    add_text(s, lx + pad, card_y + Inches(0.30), col_w - 2 * pad, Inches(0.28),
             {"text": "TARGET USER", "size": 10, "bold": True, "color": INK_500})
    add_text(s, lx + pad, card_y + Inches(0.65), col_w - 2 * pad, Inches(0.55),
             {"text": "Occasional outdoor adventurers (18–35)",
              "size": 19, "bold": True, "color": INK_900,
              "line_spacing": 1.1})

    bullet_y = card_y + Inches(1.45)
    bullets = [
        "Urban professionals & students who hike, camp or paddle 2–6 weekends a year.",
        "Cost-sensitive: a 4-day backpacking trip costs ~$1,400 in new gear they’ll use once a season.",
        "Trust-anxious: 71% have skipped a peer-to-peer rental because of safety, payment or no-show fears (M1 survey, n=42).",
        "Sustainability-conscious: 64% would prefer renting to buying if logistics were as easy.",
    ]
    add_lines(s, lx + pad, bullet_y, col_w - 2 * pad, card_h - Inches(1.85),
              [{"text": "•  " + b, "space_before": 0 if i == 0 else 8}
               for i, b in enumerate(bullets)],
              size=12, color=INK_700, line_spacing=1.30)

    # Right: problem statement on dark
    rx = lx + col_w + Inches(0.30)
    add_round(s, rx, card_y, col_w, card_h, INK_900, radius=0.04)
    add_rect(s, rx, card_y, Inches(0.10), card_h, BRAND_500)

    add_text(s, rx + pad, card_y + Inches(0.30), col_w - 2 * pad, Inches(0.28),
             {"text": "FINAL PROBLEM STATEMENT",
              "size": 10, "bold": True, "color": BRAND_400})

    # The statement — single textbox with inline runs (NO duplicate)
    add_runs(s, rx + pad, card_y + Inches(0.75), col_w - 2 * pad,
             card_h - Inches(0.95),
             [
                 {"text": "Occasional outdoor adventurers (18–35)",
                  "color": WHITE, "bold": True, "size": 17},
                 {"text": " need ", "color": INK_300, "size": 17},
                 {"text": "a trustworthy way to access quality outdoor gear "
                          "for short trips",
                  "color": WHITE, "bold": True, "size": 17},
                 {"text": " because ", "color": INK_300, "size": 17},
                 {"text": "buying gear is expensive and most kit sits idle "
                          "90% of the year",
                  "color": WHITE, "bold": True, "size": 17},
                 {"text": ", which could be met by ", "color": INK_300,
                  "size": 17},
                 {"text": "a verified, hyper-local peer-to-peer marketplace "
                          "with built-in escrow and damage protection.",
                  "color": BRAND_400, "bold": True, "size": 17},
             ],
             line_spacing=1.32)


# =========================================================
# SLIDE 2 — Key Needs, Functions and Design Priorities
# =========================================================
def build_slide_2(page_total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, INK_50)
    slide_chrome(s, page_label="2", page_total=page_total)

    add_text(s, Inches(0.85), Inches(1.05), Inches(8), Inches(0.35),
             {"text": "02 — DESIGN REQUIREMENTS",
              "size": 11, "bold": True, "color": BRAND_600})
    add_text(s, Inches(0.85), Inches(1.40), Inches(11.6), Inches(0.7),
             {"text": "What every design decision had to satisfy.",
              "size": 30, "bold": True, "color": INK_900,
              "line_spacing": 1.1})

    cw = Inches(3.95)
    ch = Inches(4.55)
    cy = Inches(2.40)
    gap = Inches(0.20)

    # ---- Column 1: Customer needs ----
    cx1 = Inches(0.85)
    add_round(s, cx1, cy, cw, ch, WHITE, radius=0.04, line=INK_200)
    add_text(s, cx1 + Inches(0.30), cy + Inches(0.30), cw - Inches(0.6),
             Inches(0.28),
             {"text": "▎ CUSTOMER NEEDS (M2)",
              "size": 10, "bold": True, "color": BRAND_600})
    add_text(s, cx1 + Inches(0.30), cy + Inches(0.65), cw - Inches(0.6),
             Inches(0.4),
             {"text": "Top 5 of 20 user stories",
              "size": 13, "bold": True, "color": INK_900})

    needs = [
        ("1", "Verified counterparty", "Both sides ID-checked"),
        ("2", "Local pickup",          "Map-based proximity search"),
        ("3", "Secure payment",        "Escrow, no cash hand-off"),
        ("4", "Damage protection",     "Up to $500 per rental"),
        ("5", "Honest condition",      "Mandatory photos + grade"),
    ]
    list_y = cy + Inches(1.30)
    row_h  = Inches(0.60)
    for i, (n, t, sub) in enumerate(needs):
        ny = list_y + row_h * i
        add_round(s, cx1 + Inches(0.30), ny + Inches(0.06),
                  Inches(0.32), Inches(0.32), BRAND_50, radius=0.5)
        add_text(s, cx1 + Inches(0.30), ny + Inches(0.06),
                 Inches(0.32), Inches(0.32),
                 {"text": n, "size": 11, "bold": True, "color": BRAND_600,
                  "align": PP_ALIGN.CENTER},
                 anchor=MSO_ANCHOR.MIDDLE)
        add_lines(s, cx1 + Inches(0.72), ny, cw - Inches(0.95), row_h,
                  [{"text": t, "size": 12, "bold": True, "color": INK_900},
                   {"text": sub, "size": 10, "color": INK_500,
                    "space_before": 2}],
                  line_spacing=1.2)

    # ---- Column 2: Functions ----
    cx2 = cx1 + cw + gap
    add_round(s, cx2, cy, cw, ch, WHITE, radius=0.04, line=INK_200)
    add_text(s, cx2 + Inches(0.30), cy + Inches(0.30), cw - Inches(0.6),
             Inches(0.28),
             {"text": "▎ PRIME FUNCTION + SUB-FUNCTIONS",
              "size": 10, "bold": True, "color": BRAND_600})
    add_text(s, cx2 + Inches(0.30), cy + Inches(0.65), cw - Inches(0.6),
             Inches(0.4),
             {"text": "Functional decomposition (M2)",
              "size": 13, "bold": True, "color": INK_900})

    add_round(s, cx2 + Inches(0.30), cy + Inches(1.30),
              cw - Inches(0.6), Inches(0.55),
              INK_900, radius=0.3)
    add_text(s, cx2 + Inches(0.30), cy + Inches(1.30),
             cw - Inches(0.6), Inches(0.55),
             {"text": "Enable trusted P2P rental of outdoor gear",
              "size": 11, "bold": True, "color": WHITE,
              "align": PP_ALIGN.CENTER},
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, cx2 + Inches(0.30), cy + Inches(1.92),
             cw - Inches(0.6), Inches(0.25),
             {"text": "▼", "size": 11, "color": INK_300,
              "align": PP_ALIGN.CENTER})

    subs = [
        ("Discover",     "Search · filter · map"),
        ("Transact",     "Book · pay · escrow"),
        ("Communicate",  "Message · notify"),
        ("Build trust",  "Verify · review · rate"),
    ]
    sub_y = cy + Inches(2.30)
    sub_h = Inches(0.42)
    for i, (k, v) in enumerate(subs):
        sy = sub_y + (sub_h + Inches(0.08)) * i
        add_round(s, cx2 + Inches(0.30), sy, cw - Inches(0.6), sub_h,
                  BRAND_50, radius=0.3)
        add_text(s, cx2 + Inches(0.45), sy, Inches(1.5), sub_h,
                 {"text": k, "size": 11, "bold": True, "color": BRAND_700},
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, cx2 + Inches(1.85), sy, cw - Inches(2.15), sub_h,
                 {"text": v, "size": 10, "color": INK_700},
                 anchor=MSO_ANCHOR.MIDDLE)

    # ---- Column 3: HoQ priorities ----
    cx3 = cx2 + cw + gap
    add_round(s, cx3, cy, cw, ch, INK_900, radius=0.04)
    add_text(s, cx3 + Inches(0.30), cy + Inches(0.30), cw - Inches(0.6),
             Inches(0.28),
             {"text": "▎ TOP HoQ TECHNICAL PRIORITIES",
              "size": 10, "bold": True, "color": BRAND_400})
    add_text(s, cx3 + Inches(0.30), cy + Inches(0.65), cw - Inches(0.6),
             Inches(0.4),
             {"text": "Highest weighted from M2 House of Quality",
              "size": 12, "bold": True, "color": WHITE,
              "line_spacing": 1.2})

    techs = [
        ("Trust score engine",      9.4, "ID + email + phone + history"),
        ("Geospatial search latency", 8.8, "PostGIS, p95 < 600 ms"),
        ("Stripe Connect escrow",    8.5, "Capture only on owner approval"),
    ]
    tech_y = cy + Inches(1.40)
    block_h = Inches(0.95)
    for i, (k, score, desc) in enumerate(techs):
        ny = tech_y + block_h * i
        add_text(s, cx3 + Inches(0.30), ny, cw - Inches(0.6), Inches(0.30),
                 {"text": k, "size": 13, "bold": True, "color": WHITE})
        add_text(s, cx3 + Inches(0.30), ny + Inches(0.32),
                 cw - Inches(0.6), Inches(0.25),
                 {"text": desc, "size": 9.5, "color": INK_300})
        bar_y = ny + Inches(0.65)
        bar_w = cw - Inches(0.6)
        add_round(s, cx3 + Inches(0.30), bar_y, bar_w, Inches(0.10),
                  INK_800, radius=0.5)
        add_round(s, cx3 + Inches(0.30), bar_y, bar_w * (score / 10.0),
                  Inches(0.10), BRAND_500, radius=0.5)
        add_text(s, cx3 + cw - Inches(0.85), bar_y - Inches(0.02),
                 Inches(0.55), Inches(0.20),
                 {"text": f"{score:.1f}", "size": 10, "bold": True,
                  "color": BRAND_400, "align": PP_ALIGN.RIGHT},
                 anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 3 — Concept Selection & Design Evolution (M3 → M4)
# =========================================================
def build_slide_3(page_total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, INK_50)
    slide_chrome(s, page_label="3", page_total=page_total)

    add_text(s, Inches(0.85), Inches(1.05), Inches(8), Inches(0.35),
             {"text": "03 — CONCEPT SELECTION & DESIGN EVOLUTION",
              "size": 11, "bold": True, "color": BRAND_600})
    add_text(s, Inches(0.85), Inches(1.40), Inches(11.6), Inches(0.7),
             {"text": "Why Concept C, and what changed between M3 and M4.",
              "size": 26, "bold": True, "color": INK_900,
              "line_spacing": 1.1})

    # ---- Pugh chart row ----
    pcy = Inches(2.30)
    pch = Inches(1.95)
    add_round(s, Inches(0.85), pcy, Inches(12.0), pch, WHITE, radius=0.04,
              line=INK_200)
    add_text(s, Inches(1.10), pcy + Inches(0.18), Inches(11.5), Inches(0.30),
             {"text": "PUGH CHART RESULT (M3) — net weighted score across 10 criteria",
              "size": 10, "bold": True, "color": INK_500})

    concepts = [
        ("A · Classifieds",   0,  INK_300, False),
        ("B · Full-Managed",  0,  INK_300, False),
        ("C · Trust-First",   30, BRAND_500, True),
        ("D · Subscription", -3,  ROSE_500, False),
        ("E · AI Matching",   3,  AMBER_500, False),
    ]
    bar_x = Inches(1.10)
    bar_w = Inches(11.5)
    bar_y = pcy + Inches(0.62)
    bar_h = Inches(0.46)
    seg_w = bar_w / len(concepts)
    max_abs = 30
    for i, (name, score, color, winner) in enumerate(concepts):
        sx = bar_x + seg_w * i + Inches(0.08)
        sw_ = seg_w - Inches(0.16)
        # Track
        add_round(s, sx, bar_y, sw_, bar_h, INK_100, radius=0.3)
        if score >= 0:
            fw = sw_ * (score / max_abs)
            add_round(s, sx, bar_y, fw, bar_h, color, radius=0.3)
        else:
            fw = sw_ * (abs(score) / max_abs)
            add_round(s, sx + sw_ - fw, bar_y, fw, bar_h, color, radius=0.3)
        add_runs(s, sx, bar_y + bar_h + Inches(0.06), sw_, Inches(0.30),
                 [{"text": name + "   ", "size": 11,
                   "color": INK_900 if winner else INK_700,
                   "bold": winner},
                  {"text": f"net {score:+d}", "size": 11, "bold": True,
                   "color": color if (winner or score < 0) else INK_500}],
                 align=PP_ALIGN.CENTER)

    # ---- Reason callout: WHY Concept C was selected ----
    reason_y = pcy + Inches(1.45)
    add_runs(s, Inches(1.10), reason_y, Inches(11.5), Inches(0.40),
             [{"text": "→ Why C? ", "size": 11, "bold": True,
               "color": BRAND_600},
              {"text": "Only concept with ", "size": 11, "color": INK_700},
              {"text": "zero negatives across all 10 criteria",
               "size": 11, "bold": True, "color": INK_900},
              {"text": " — combines the trust infrastructure of B (full-managed) with the asset-light scalability of A (classifieds), and uniquely satisfies the M1 trust-deficit insight.",
               "size": 11, "color": INK_700}],
             line_spacing=1.30,
             anchor=MSO_ANCHOR.MIDDLE)

    # ---- M3 → M4 modifications table ----
    table_y = pcy + pch + Inches(0.30)
    table_w = Inches(12.0)
    cols = [
        ("M3 issue / feedback",                 4.30),
        ("Modification in M4",                  4.70),
        ("Reason for change (DfX method)",      3.00),
    ]

    # Header
    hdr_h = Inches(0.42)
    add_round(s, Inches(0.85), table_y, table_w, hdr_h, INK_900, radius=0.04)
    cx = Inches(0.85)
    for header, ww in cols:
        add_text(s, cx + Inches(0.25), table_y, Inches(ww), hdr_h,
                 {"text": header, "size": 10.5, "bold": True,
                  "color": BRAND_400},
                 anchor=MSO_ANCHOR.MIDDLE)
        cx += Inches(ww)

    # Data
    rows = [
        ("Booking page was one wall of fields → high drop-off risk",
         "Split into 3-step Stepper: Review → Pay → Confirm",
         "DfU — reduces cognitive load (Miller 7±2)"),
        ("Pickup address visible before owner approval",
         "Address only revealed after escrow capture",
         "DfU — privacy + trust handshake"),
        ("Listing card lacked any verification cue",
         "Added “✓ Verified Owner” chip on every card",
         "DfU — surfaces trust at scan-time"),
        ("3.8 : 1 contrast on hero overlay; 32 px tap targets",
         "Lifted contrast to 4.7 : 1; all CTAs ≥ 44 × 44 px",
         "DfA — WCAG 2.1 AA compliance"),
    ]
    row_h = Inches(0.55)
    for i, row in enumerate(rows):
        ry = table_y + hdr_h + row_h * i
        bg = WHITE if i % 2 == 0 else INK_50
        add_rect(s, Inches(0.85), ry, table_w, row_h, bg)
        cx = Inches(0.85)
        for j, (text, (_, ww)) in enumerate(zip(row, cols)):
            color = INK_900 if j == 1 else INK_700
            add_text(s, cx + Inches(0.25), ry, Inches(ww), row_h,
                     {"text": text, "size": 10, "bold": (j == 1),
                      "color": color, "line_spacing": 1.15},
                     anchor=MSO_ANCHOR.MIDDLE)
            cx += Inches(ww)

    # Outline
    out = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.85), table_y,
                             table_w, hdr_h + row_h * len(rows))
    out.fill.background()
    out.line.color.rgb = INK_200
    out.line.width = Pt(0.75)
    out.shadow.inherit = False


# =========================================================
# SLIDE 4 — Final Embodiment Design
# =========================================================
def build_slide_4(page_total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, INK_50)
    slide_chrome(s, page_label="4", page_total=page_total)

    add_text(s, Inches(0.85), Inches(1.05), Inches(8), Inches(0.35),
             {"text": "04 — FINAL EMBODIMENT DESIGN",
              "size": 11, "bold": True, "color": BRAND_600})
    add_text(s, Inches(0.85), Inches(1.40), Inches(8.5), Inches(0.55),
             {"text": "EcoRent v1 — built end-to-end in Next.js 14.",
              "size": 24, "bold": True, "color": INK_900,
              "line_spacing": 1.1})
    add_text(s, Inches(0.85), Inches(2.00), Inches(8.5), Inches(0.30),
             {"text": "Live demo deployed on Vercel. Below: four core screens, the final user flow and the architecture.",
              "size": 10.5, "color": INK_500})

    # ---- Live demo URL callout (top-right) ----
    demo_x = Inches(9.55)
    demo_y = Inches(1.30)
    demo_w = Inches(3.30)
    demo_h = Inches(0.95)
    add_round(s, demo_x, demo_y, demo_w, demo_h, INK_900, radius=0.06,
              line=BRAND_500, line_w=1.25)
    add_text(s, demo_x + Inches(0.20), demo_y + Inches(0.14),
             demo_w - Inches(0.40), Inches(0.24),
             {"text": "▶  TRY THE LIVE DEMO",
              "size": 9.5, "bold": True, "color": BRAND_400})
    add_text(s, demo_x + Inches(0.20), demo_y + Inches(0.40),
             demo_w - Inches(0.40), Inches(0.40),
             {"text": "pdd-azure.vercel.app",
              "size": 16, "bold": True, "color": WHITE,
              "line_spacing": 1.0})

    # ---- Top: 4 screen thumbnails ----
    sy = Inches(2.45)
    th = Inches(2.30)
    pad = Inches(0.16)
    tw = (Inches(12.0) - 3 * pad) / 4
    screens = [
        ("landing.png",   "Landing"),
        ("search.png",    "Search · Map"),
        ("listing.png",   "Listing Detail"),
        ("dashboard.png", "Dashboard"),
    ]
    for i, (img, label) in enumerate(screens):
        x = Inches(0.85) + (tw + pad) * i
        add_round(s, x, sy, tw, th, INK_900, radius=0.04, line=INK_200)
        add_text(s, x + Inches(0.20), sy + Inches(0.10),
                 tw - Inches(0.4), Inches(0.30),
                 {"text": label, "size": 11, "bold": True,
                  "color": BRAND_400})
        img_path = ASSETS / img
        if img_path.exists():
            try:
                s.shapes.add_picture(
                    str(img_path),
                    x + Inches(0.12),
                    sy + Inches(0.50),
                    width=tw - Inches(0.24),
                    height=th - Inches(0.62),
                )
            except Exception:
                pass

    # ---- Bottom: user flow + arch + features ----
    by = Inches(5.00)
    bh = Inches(2.10)

    # Left card: user flow + features stacked
    flow_w = Inches(7.65)
    add_round(s, Inches(0.85), by, flow_w, bh, WHITE, radius=0.04,
              line=INK_200)
    add_text(s, Inches(1.05), by + Inches(0.18),
             flow_w - Inches(0.4), Inches(0.25),
             {"text": "FINAL USER FLOW",
              "size": 10, "bold": True, "color": BRAND_600})

    flow = ["Land", "Search", "Listing", "Book", "Confirm", "Pickup", "Review"]
    fy = by + Inches(0.55)
    pill_h = Inches(0.50)
    arrow_w = Inches(0.16)
    n = len(flow)
    avail = flow_w - Inches(0.40) - arrow_w * (n - 1)
    pill_w = avail / n
    fx = Inches(1.05)
    for i, step in enumerate(flow):
        px = fx + (pill_w + arrow_w) * i
        add_round(s, px, fy, pill_w, pill_h, BRAND_50, radius=0.5,
                  line=BRAND_500, line_w=1.0)
        add_text(s, px, fy, pill_w, pill_h,
                 {"text": step, "size": 10, "bold": True, "color": BRAND_700,
                  "align": PP_ALIGN.CENTER},
                 anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            add_text(s, px + pill_w, fy, arrow_w, pill_h,
                     {"text": "→", "size": 12, "bold": True,
                      "color": INK_300, "align": PP_ALIGN.CENTER},
                     anchor=MSO_ANCHOR.MIDDLE)

    # Features grid (4 mini cards)
    feats = [
        ("Trust handshake", "ID + email + phone verification chips"),
        ("Escrow booking",  "Stripe Connect — captured on approval"),
        ("Geospatial map",  "PostGIS · grid ↔ map toggle"),
        ("Sustainability",  "Per-rental CO₂ savings on receipts"),
    ]
    fy2 = by + Inches(1.20)
    fh = Inches(0.78)
    fpad = Inches(0.12)
    fw_ = (flow_w - Inches(0.40) - fpad * 3) / 4
    for i, (k, v) in enumerate(feats):
        px = Inches(1.05) + (fw_ + fpad) * i
        add_round(s, px, fy2, fw_, fh, INK_50, radius=0.06, line=INK_200)
        add_text(s, px + Inches(0.12), fy2 + Inches(0.10),
                 fw_ - Inches(0.24), Inches(0.30),
                 {"text": k, "size": 10, "bold": True, "color": INK_900})
        add_text(s, px + Inches(0.12), fy2 + Inches(0.40),
                 fw_ - Inches(0.24), Inches(0.34),
                 {"text": v, "size": 8.5, "color": INK_500,
                  "line_spacing": 1.20})

    # Right card: architecture
    ax = Inches(0.85) + flow_w + Inches(0.25)
    aw = Inches(12.0) - flow_w - Inches(0.25)
    add_round(s, ax, by, aw, bh, INK_900, radius=0.04)
    add_text(s, ax + Inches(0.20), by + Inches(0.18),
             aw - Inches(0.4), Inches(0.25),
             {"text": "PRODUCT ARCHITECTURE",
              "size": 10, "bold": True, "color": BRAND_400})

    layers = [
        ("Frontend",  "Next.js 14 · Tailwind CSS · React 18"),
        ("API",       "Django REST · OpenAPI · webhooks"),
        ("Data",      "PostgreSQL + PostGIS · Redis"),
        ("Services",  "Stripe Connect · Mapbox · S3 + CloudFront"),
    ]
    ly = by + Inches(0.55)
    line_h = Inches(0.36)
    for i, (k, v) in enumerate(layers):
        py_ = ly + line_h * i
        add_text(s, ax + Inches(0.20), py_, Inches(1.10), line_h,
                 {"text": k, "size": 10, "bold": True, "color": BRAND_400},
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, ax + Inches(1.40), py_, aw - Inches(1.60), line_h,
                 {"text": v, "size": 9.5, "color": INK_300,
                  "line_spacing": 1.2},
                 anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 5 — Conclusion
# =========================================================
def build_slide_5(page_total):
    s = prs.slides.add_slide(BLANK)
    add_rect(s, Emu(0), Emu(0), SW, SH, INK_900)
    slide_chrome(s, page_label="5", page_total=page_total, dark=True)

    add_text(s, Inches(0.85), Inches(1.05), Inches(8), Inches(0.35),
             {"text": "05 — CONCLUSION",
              "size": 11, "bold": True, "color": BRAND_400})
    add_text(s, Inches(0.85), Inches(1.40), Inches(11.6), Inches(0.55),
             {"text": "Our final design addresses the user need by…",
              "size": 22, "color": INK_300, "line_spacing": 1.1})

    # Quote card
    qy = Inches(2.30)
    qh = Inches(2.55)
    add_round(s, Inches(0.85), qy, Inches(12.0), qh, DARK_PANEL, radius=0.04)
    add_rect(s, Inches(0.85), qy, Inches(0.10), qh, BRAND_500)

    add_runs(s, Inches(1.20), qy + Inches(0.50),
             Inches(11.4), qh - Inches(1.0),
             [
                 {"text": "Combining a hyper-local peer-to-peer marketplace with ",
                  "color": INK_300, "size": 21},
                 {"text": "verified profiles, escrow payments and a damage-protection plan",
                  "color": WHITE, "bold": True, "size": 21},
                 {"text": " — making it both safe and affordable for occasional adventurers to access quality gear, while owners earn from kit that would otherwise sit idle.",
                  "color": INK_300, "size": 21},
             ],
             line_spacing=1.32)

    # Future scope row
    fy = Inches(5.20)
    fh = Inches(1.50)
    cards = [
        ("CLOSED BETA",
         "50 verified owners · Boulder, CO · summer 2026"),
        ("VALIDATE",
         "Test against M2 NFRs · iterate dispute resolution"),
        ("EXPAND",
         "Second city · Stripe Identity-only countries"),
    ]
    cw = (Inches(12.0) - Inches(0.4)) / 3
    for i, (k, v) in enumerate(cards):
        cx = Inches(0.85) + (cw + Inches(0.20)) * i
        add_round(s, cx, fy, cw, fh, DARK_PANEL, radius=0.04)
        add_text(s, cx + Inches(0.30), fy + Inches(0.22),
                 cw - Inches(0.6), Inches(0.30),
                 {"text": k, "size": 10, "bold": True, "color": BRAND_400})
        add_text(s, cx + Inches(0.30), fy + Inches(0.62),
                 cw - Inches(0.6), Inches(0.85),
                 {"text": v, "size": 12, "color": INK_100,
                  "line_spacing": 1.30})

    # Sign-off
    add_text(s, Inches(0.85), Inches(7.00), Inches(12.0), Inches(0.30),
             {"text": "Thank you. — Ashish K. Cheruku · 2021B3AB1141P",
              "size": 11, "color": INK_300, "align": PP_ALIGN.CENTER})


# =========================================================
# Build
# =========================================================
PAGE_TOTAL = 5
build_title()
build_slide_1(PAGE_TOTAL)
build_slide_2(PAGE_TOTAL)
build_slide_3(PAGE_TOTAL)
build_slide_4(PAGE_TOTAL)
build_slide_5(PAGE_TOTAL)

out = ASSETS.parent.parent / "EcoRent_M4_Final_Presentation.pptx"
prs.save(out)
print(f"saved {out}  ({out.stat().st_size/1024:.0f} KB, "
      f"{len(prs.slides)} slides)")
